from pptx import Presentation
import os
from dotenv import load_dotenv

load_dotenv()

template = os.getenv("TEMPLATE_ESR_A")
print(template)
presentation = Presentation(template)


def list_textboxes(presentation, slide_number):
    slide = presentation.slides[slide_number - 1]
    textboxes = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            textboxes.append(shape.text)
    return textboxes


# print(f"List of textboxes in the first slide: {list_textboxes(presentation, 4)}")


for idx, text in enumerate(list_textboxes(presentation, 1), 1):
    print(f"TextBox {idx}: {text}")


def update_textbox(presentation, slide_number, textbox_number, new_text):
    slide = presentation.slides[slide_number - 1]
    textboxes = [shape for shape in slide.shapes if hasattr(shape, "text")]
    if 0 < textbox_number <= len(textboxes):
        textbox = textboxes[textbox_number - 1]
        textbox.text = new_text
        print(
            f"Updated TextBox {textbox_number} on Slide {slide_number} to: {new_text}"
        )
    else:
        print(f"Textbox number {textbox_number} does not exist on Slide {slide_number}")


update_textbox(presentation, 1, 2, "ORBIT 2.0")


def save_presentation(presentation, output_path):
    presentation.save(output_path)
    print(f"Presentation saved to {output_path}")


save_presentation(presentation, "./temp/ESR-A_ORBIT20_20250415.pptx")
print("Presentation saved")
