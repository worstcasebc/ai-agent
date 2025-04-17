from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Bilderpfade (Bitte durch echte Pfade zu Bildern ersetzen)
IMAGE_PATHS = {
    "intro": "./images/llm_intro.jpg",
    "use_cases": "./images/use_cases.jpg",
    "architecture": "./images/architecture.jpg",
    "future": "./images/future.jpg",
}

# Präsentation erstellen
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Titel-Folie
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Die Macht von Large Language Models (LLMs)"
subtitle = slide.placeholders[1]
subtitle.text = "Einführung für Studierende"

# Folie 2 – Einführung in LLMs
slide = prs.slides.add_slide(prs.slide_layouts[6])
left = Inches(0.5)
top = Inches(0.5)
width = Inches(6)
height = Inches(6)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "Was sind LLMs?"
tf.paragraphs[0].font.size = Pt(36)

p = tf.add_paragraph()
p.text = (
    "Large Language Models sind KI-Modelle, die mit riesigen Textmengen trainiert wurden, "
    "um Sprache zu verstehen und zu generieren."
)
p.font.size = Pt(20)

# Bild
if os.path.exists(IMAGE_PATHS["intro"]):
    slide.shapes.add_picture(
        IMAGE_PATHS["intro"], Inches(7), Inches(1), width=Inches(5.5)
    )

# Folie 3 – Anwendungsfälle
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture(
    IMAGE_PATHS["use_cases"], Inches(0.5), Inches(0.5), width=Inches(5.5)
)

textbox = slide.shapes.add_textbox(Inches(6.3), Inches(0.5), Inches(6.5), Inches(6))
tf = textbox.text_frame
tf.text = "Wo werden LLMs genutzt?"
tf.paragraphs[0].font.size = Pt(30)

for case in [
    "Texterstellung (z. B. Chatbots, Artikel)",
    "Code-Generierung (z. B. GitHub Copilot)",
    "Sprachübersetzung",
    "Datenanalyse",
]:
    p = tf.add_paragraph()
    p.text = f"• {case}"
    p.font.size = Pt(20)

# Folie 4 – Architektur
slide = prs.slides.add_slide(prs.slide_layouts[6])
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(6))
tf = textbox.text_frame
tf.text = "Wie funktionieren LLMs technisch?"
tf.paragraphs[0].font.size = Pt(30)

p = tf.add_paragraph()
p.text = "Transformator-Architektur mit Self-Attention und Deep Learning auf riesigen Datenmengen."
p.font.size = Pt(20)

if os.path.exists(IMAGE_PATHS["architecture"]):
    slide.shapes.add_picture(
        IMAGE_PATHS["architecture"], Inches(6.5), Inches(1), width=Inches(6)
    )

# Folie 5 – Zukunftsaussichten
slide = prs.slides.add_slide(prs.slide_layouts[6])
if os.path.exists(IMAGE_PATHS["future"]):
    slide.shapes.add_picture(
        IMAGE_PATHS["future"], Inches(0.5), Inches(0.5), width=Inches(6)
    )

textbox = slide.shapes.add_textbox(Inches(6.7), Inches(0.5), Inches(6.5), Inches(6.5))
tf = textbox.text_frame
tf.text = "Wohin geht die Reise?"
tf.paragraphs[0].font.size = Pt(28)

for trend in [
    "Multimodale KI",
    "Personalisierte Assistenzsysteme",
    "Ethische Herausforderungen",
    "Autonome Agenten",
]:
    p = tf.add_paragraph()
    p.text = f"- {trend}"
    p.font.size = Pt(20)

# Präsentation speichern
prs.save("./temp/llm_praesentation.pptx")
print("PowerPoint-Präsentation erfolgreich erstellt.")
